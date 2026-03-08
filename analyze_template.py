#!/usr/bin/env python3
"""Analyze the Atlan Beta User Training PPTX for style extraction."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

def safe_color(color_obj):
    try:
        if color_obj and color_obj.type is not None:
            return f"#{color_obj.rgb}"
    except:
        try:
            return f"theme={color_obj.theme_color}"
        except:
            pass
    return None

prs = Presentation("/Users/dattu.marneni/Desktop/Sun Aug 31 16-53-08 2025/Atlan_Beta_User_Training.pptx")

print(f"Slide width: {prs.slide_width / 914400:.2f} inches")
print(f"Slide height: {prs.slide_height / 914400:.2f} inches")
print(f"Total slides: {len(prs.slides)}")

for i, layout in enumerate(prs.slide_layouts):
    print(f"  Layout {i}: {layout.name}")

for slide_idx, slide in enumerate(prs.slides):
    print(f"\n{'='*80}")
    print(f"SLIDE {slide_idx + 1}")
    print(f"{'='*80}")
    
    try:
        bg = slide.background
        if bg.fill.type is not None:
            fc = safe_color(bg.fill.fore_color)
            print(f"  Background: type={bg.fill.type}, color={fc}")
    except:
        print("  Background: (complex)")
    
    for shape_idx, shape in enumerate(slide.shapes):
        print(f"\n  Shape {shape_idx}: type={shape.shape_type} | name='{shape.name}'")
        print(f"    Pos: ({shape.left/914400:.2f}in, {shape.top/914400:.2f}in) Size: ({shape.width/914400:.2f}in x {shape.height/914400:.2f}in)")
        
        try:
            fill = shape.fill
            if fill.type is not None:
                fc = safe_color(fill.fore_color)
                print(f"    Fill: type={fill.type}, color={fc}")
        except:
            pass
        
        if shape.has_text_frame:
            tf = shape.text_frame
            for p_idx, para in enumerate(tf.paragraphs):
                if para.text.strip():
                    text_preview = para.text[:80]
                    info = []
                    
                    for run in para.runs:
                        if run.font.size:
                            info.append(f"size={run.font.size.pt:.0f}pt")
                        if run.font.bold:
                            info.append("bold")
                        if run.font.italic:
                            info.append("italic")
                        c = safe_color(run.font.color)
                        if c:
                            info.append(f"color={c}")
                        if run.font.name:
                            info.append(f"font={run.font.name}")
                        break
                    
                    if para.alignment:
                        info.append(f"align={para.alignment}")
                    
                    info_str = " | ".join(info) if info else "default"
                    print(f"    P{p_idx}: \"{text_preview}\" [{info_str}]")
        
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                print(f"    [IMAGE] {shape.image.content_type}")
            except:
                print(f"    [IMAGE]")
        
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"    [GROUP] {len(shape.shapes)} sub-shapes")
            for ss_idx, ss in enumerate(shape.shapes):
                print(f"      Sub {ss_idx}: type={ss.shape_type} name='{ss.name}'")
                if ss.has_text_frame:
                    for p in ss.text_frame.paragraphs:
                        if p.text.strip():
                            rinfo = []
                            for r in p.runs:
                                if r.font.size: rinfo.append(f"size={r.font.size.pt:.0f}pt")
                                if r.font.bold: rinfo.append("bold")
                                c = safe_color(r.font.color)
                                if c: rinfo.append(f"color={c}")
                                if r.font.name: rinfo.append(f"font={r.font.name}")
                                break
                            print(f"        \"{p.text[:60]}\" [{', '.join(rinfo)}]")
