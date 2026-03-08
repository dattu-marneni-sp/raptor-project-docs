#!/usr/bin/env python3
"""
Generate Raptor presentations combining:
- Rich visual content (cards, metrics, color-coded lanes) from the first version
- SailPoint branded template backgrounds from Atlan_Beta_User_Training.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE = "/Users/dattu.marneni/Desktop/Sun Aug 31 16-53-08 2025/Atlan_Beta_User_Training.pptx"

# SailPoint brand colors
NAVY = RGBColor(0x0D, 0x1B, 0x3E)
DARK_BLUE = RGBColor(0x14, 0x27, 0x52)
SAILPOINT_BLUE = RGBColor(0x00, 0x33, 0x66)
TEAL = RGBColor(0x00, 0xB4, 0xD8)
LIGHT_TEAL = RGBColor(0x48, 0xCA, 0xE4)
ORANGE = RGBColor(0xF7, 0x7F, 0x00)
LIGHT_ORANGE = RGBColor(0xFF, 0xA6, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
MID_GRAY = RGBColor(0x8C, 0x8C, 0xA0)
DARK_GRAY = RGBColor(0x3A, 0x3A, 0x4A)
GREEN = RGBColor(0x28, 0xA7, 0x45)
YELLOW = RGBColor(0xFF, 0xC1, 0x07)
RED = RGBColor(0xDC, 0x35, 0x45)
ACCENT_PURPLE = RGBColor(0x6F, 0x42, 0xC1)
BLUE_ACCENT = RGBColor(0x20, 0x26, 0xD2)
CARD_BG = RGBColor(0x1A, 0x2D, 0x5A)
SEMI_TRANSPARENT = RGBColor(0x10, 0x18, 0x35)

LAYOUT_TITLE = 0
LAYOUT_SECTION_DARK = 23
LAYOUT_CONTENT_DARK = 8
LAYOUT_CONTENT = 7
LAYOUT_BLANK = 25


def new_prs():
    prs = Presentation(TEMPLATE)
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId is None:
            rId = sldId.get('r:id')
        if rId:
            try:
                prs.part.drop_rel(rId)
            except:
                pass
        sldIdLst.remove(sldId)
    return prs


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Poppins"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_card(slide, left, top, width, height, title, body_lines, icon="",
             header_color=TEAL):
    card = add_shape(slide, left, top, width, height, CARD_BG, border_color=header_color, border_width=1.5)
    add_rect(slide, left, top, width, Inches(0.06), header_color)

    title_text = f"{icon}  {title}" if icon else title
    add_text_box(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.4),
                 title_text, font_size=14, font_color=header_color, bold=True)

    y_offset = top + Inches(0.55)
    for line in body_lines:
        add_text_box(slide, left + Inches(0.15), y_offset, width - Inches(0.3), Inches(0.28),
                     line, font_size=11, font_color=LIGHT_GRAY, bold=False)
        y_offset += Inches(0.25)


def add_metric_card(slide, left, top, width, height, metric_value, metric_label, accent_color=TEAL):
    add_shape(slide, left, top, width, height, CARD_BG, border_color=accent_color, border_width=2)
    add_text_box(slide, left, top + Inches(0.15), width, Inches(0.5),
                 metric_value, font_size=32, font_color=accent_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.65), width - Inches(0.2), Inches(0.4),
                 metric_label, font_size=11, font_color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.CENTER)


def make_title_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = ""
            tf = shape.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
            run.font.name = "Poppins SemiBold"
            run.font.size = Pt(40)
            run.font.color.rgb = WHITE
    return slide


def make_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION_DARK])
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = ""
            tf = shape.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
            run.font.name = "Poppins SemiBold"
            run.font.size = Pt(36)
            run.font.color.rgb = WHITE
    return slide


def make_content_slide(prs, title, dark=True):
    layout = LAYOUT_CONTENT_DARK if dark else LAYOUT_CONTENT
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = ""
            tf = shape.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
            run.font.name = "Poppins SemiBold"
            run.font.size = Pt(26)
            run.font.color.rgb = WHITE if dark else RGBColor(0x33, 0x33, 0x33)
        if shape.placeholder_format.idx in (1, 10):
            shape.text = ""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 12:
            shape.text = "© 2026 SailPoint Technologies, Inc. All rights reserved."
    return slide


# ==============================================================================
# PPT 1: EXECUTIVE OVERVIEW
# ==============================================================================
def create_exec_overview():
    prs = new_prs()

    # 1. Title
    make_title_slide(prs, "🦅  Raptor\nOne Door to Data at Scale")

    # 2. The Challenge - 5 problem cards
    slide = make_content_slide(prs, "The Challenge: Why We Need Raptor")

    problems = [
        ("💰", "Unsustainable Cost", "Materializers cost $100K+/mo\nRe-processing millions of\nrecords for small changes", ORANGE),
        ("🕸️", "Growing Complexity", "Each feature adds bespoke\npipelines & schema joins\nBrittle web of dependencies", YELLOW),
        ("🚧", "Platform Bottleneck", "Data platform is critical\npath for all teams\nSmall changes need DE work", RED),
        ("⏰", "Fragile Freshness", "\"Real-time\" = batch runs\nHours of lag before fresh\ndata reaches products", ACCENT_PURPLE),
        ("💣", "High Blast Radius", "Minor updates cascade\nthrough millions of rows\nDays of re-materialization", TEAL),
    ]

    for i, (icon, title, desc, color) in enumerate(problems):
        left = Inches(0.5) + Inches(i * 2.5)
        add_card(slide, left, Inches(1.8), Inches(2.3), Inches(2.6),
                 title, desc.split("\n"), icon=icon, header_color=color)

    add_text_box(slide, Inches(0.5), Inches(4.8), Inches(12), Inches(0.4),
                 "⚠️  Current search materialization costs exceed $150,000/month in compute alone",
                 font_size=16, font_color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

    add_shape(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(1.2),
              CARD_BG, border_color=TEAL, border_width=1)
    add_text_box(slide, Inches(2), Inches(5.6), Inches(9.3), Inches(1.0),
                 "\"Our current approach is hitting a wall — technically, financially, and organizationally.\n"
                 "Raptor is the necessary pivot — from pre-built documents to composable, on-demand data delivery.\"",
                 font_size=13, font_color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

    # 3. What is Raptor - Before/After
    slide = make_content_slide(prs, "What is Raptor?")

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7),
                 "A single external API that powers every product experience needing data,\n"
                 "delivered quickly, accurately, and cost-effectively.",
                 font_size=18, font_color=LIGHT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    add_card(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(3.2),
             "BEFORE  (Materializer-first)", [
                 "❌  Big pre-built documents per use case",
                 "❌  Expensive backfills & replays",
                 "❌  Multi-day lead time for new fields",
                 "❌  Freshness via constant re-materialization",
                 "❌  Every change = millions rebuilt",
                 "❌  High storage from data duplication",
             ], icon="🔴", header_color=RED)

    add_card(slide, Inches(7), Inches(2.6), Inches(5.5), Inches(3.2),
             "AFTER  (Raptor / GraphQL-first)", [
                 "✅  On-demand composition via federated subgraphs",
                 "✅  Real-time by design: fetch fresh state",
                 "✅  New fields in hours, not weeks",
                 "✅  Materializers become thin or optional",
                 "✅  Targeted rebuilds, not global reprocessing",
                 "✅  Single source of truth per entity",
             ], icon="🟢", header_color=GREEN)

    add_shape(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.6),
              CARD_BG, border_color=ORANGE, border_width=1)
    add_text_box(slide, Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.5),
                 "🍽️  Think of it like a restaurant: Today we cook every item just in case. "
                 "With Raptor, we prepare ingredients once and cook dishes fresh when ordered.",
                 font_size=12, font_color=LIGHT_ORANGE, bold=False, alignment=PP_ALIGN.CENTER)

    # 4. Architecture - 5 Lanes
    slide = make_content_slide(prs, "Architecture: 5 Specialized Data Lanes")

    lanes = [
        ("⚡", "Pointed\nLookups", "ScyllaDB / Aerospike", "p95 ≤ 10ms", "Entity hydration\nin milliseconds", TEAL),
        ("🔗", "Graph\nDatabase", "JanusGraph on Scylla", "p95 < 1s", "Relationship\ntraversals", GREEN),
        ("📊", "Real-Time\nAggregates", "ClickHouse", "p95 < 300ms", "Dashboards &\nanomalies", ORANGE),
        ("🔍", "Search &\nDiscovery", "OpenSearch / ES", "p95 < 200ms", "Full-text search\n& autocomplete", ACCENT_PURPLE),
        ("📚", "Historical\n& Batch", "Snowflake", "Hours", "Reporting,\ncompliance, ML", MID_GRAY),
    ]

    for i, (icon, name, backend, slo, purpose, color) in enumerate(lanes):
        left = Inches(0.5) + Inches(i * 2.5)
        add_shape(slide, left, Inches(1.8), Inches(2.3), Inches(3.0),
                  CARD_BG, border_color=color, border_width=2)
        add_rect(slide, left, Inches(1.8), Inches(2.3), Inches(0.06), color)

        add_text_box(slide, left, Inches(1.95), Inches(2.3), Inches(0.4),
                     icon, font_size=28, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.35), Inches(2.3), Inches(0.5),
                     name, font_size=13, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.85), Inches(2.3), Inches(0.25),
                     backend, font_size=9, font_color=color, bold=False, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(3.15), Inches(2.3), Inches(0.35),
                     slo, font_size=20, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.1), Inches(3.55), Inches(2.1), Inches(0.5),
                     purpose, font_size=9, font_color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

    add_shape(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.7),
              CARD_BG, border_color=TEAL, border_width=2)
    add_text_box(slide, Inches(0.5), Inches(5.15), Inches(12.3), Inches(0.6),
                 "🌐  Unified GraphQL Gateway  (Apollo Router + Federation 2)  →  Single External Schema",
                 font_size=16, font_color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.4),
                 "🤖  AI-Ready via MCP: Schema & metadata auto-published for AI agents and copilots",
                 font_size=13, font_color=LIGHT_ORANGE, bold=False, alignment=PP_ALIGN.CENTER)

    # 5. Target Outcomes - Metric Cards + Impact Cards
    slide = make_content_slide(prs, "Target Outcomes & Business Impact")

    metrics = [
        ("30-50%", "Cost Reduction\nvs. Snowflake", TEAL),
        ("< 2 days", "New Entity\nto GraphQL", GREEN),
        ("≥ 3", "Materializers\nRetired", ORANGE),
        ("99.99%", "Availability\nTarget", ACCENT_PURPLE),
        ("8/10", "Developer\nSatisfaction", LIGHT_TEAL),
    ]
    for i, (val, label, color) in enumerate(metrics):
        left = Inches(0.5) + Inches(i * 2.5)
        add_metric_card(slide, left, Inches(1.8), Inches(2.3), Inches(1.4), val, label, color)

    impacts = [
        ("🚀", "Faster Delivery", "No custom pipelines needed.\nDays instead of weeks."),
        ("⏱️", "Real-Time Experiences", "JIT access, live dashboards,\ndynamic policy evaluation."),
        ("💰", "Lower Cost", "Reduce data duplication,\nfewer breakages."),
        ("🤖", "AI-Ready Foundation", "GraphQL + MCP enables\nAI copilot integration."),
    ]
    for i, (icon, title, desc) in enumerate(impacts):
        left = Inches(0.5) + Inches(i * 3.1)
        add_card(slide, left, Inches(3.6), Inches(2.9), Inches(1.7),
                 title, desc.split("\n"), icon=icon, header_color=TEAL)

    add_text_box(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4),
                 "\"Raptor gives product teams a single, flexible, real-time source of truth for every data experience.\"",
                 font_size=14, font_color=LIGHT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    # 6. Next Steps
    slide = make_content_slide(prs, "⚡ Immediate Next Steps — Raptor Transition")

    actions = [
        ("1️⃣", "🔍  Assess current POC state", "Review deployed environments: Apollo, ScyllaDB, JanusGraph"),
        ("2️⃣", "👑  Clarify ownership", "Establish new drivers for each DACI lane"),
        ("3️⃣", "🎯  Address DPINTAKE-33", "Delivery date + incremental plan for Identity Graph"),
        ("4️⃣", "🔓  Unblock backlog", "Assign subgraph stories (5980/5981/5982)"),
        ("5️⃣", "🦅  Rename to Raptor", "Update Confluence, Jira labels, Slack channels"),
        ("6️⃣", "📅  Establish cadence", "Weekly \"Raptor Updates\" posts"),
        ("7️⃣", "🤝  Schedule kickoff sync", "Forge, Keel, DIPO, Double Zero leads"),
    ]
    for i, (num, action, detail) in enumerate(actions):
        top = Inches(1.7) + Inches(i * 0.72)
        add_rect(slide, Inches(0.8), top, Inches(0.08), Inches(0.55), TEAL if i % 2 == 0 else ORANGE)
        add_text_box(slide, Inches(1.1), top + Inches(0.02), Inches(5), Inches(0.3),
                     f"{num}  {action}", font_size=15, font_color=WHITE, bold=True)
        add_text_box(slide, Inches(1.1), top + Inches(0.32), Inches(11), Inches(0.25),
                     detail, font_size=12, font_color=LIGHT_GRAY, bold=False)

    # 7. Q&A
    make_section_slide(prs, "Q&A\n\n📧 dattu.marneni@sailpoint.com\n📋 DPINTAKE-33\n💬 #raptor-data-platform")

    prs.save("Raptor_Executive_Overview.pptx")
    print("✅  Created: Raptor_Executive_Overview.pptx")


# ==============================================================================
# PPT 2: TECHNICAL DEEP DIVE
# ==============================================================================
def create_tech_deep_dive():
    prs = new_prs()

    make_title_slide(prs, "🦅  Raptor\nTechnical Deep Dive")

    # GraphQL Federation Flow
    make_section_slide(prs, "GraphQL Federation\nHow the Unified API Works")

    slide = make_content_slide(prs, "Data Flow: Source → Lane → GraphQL → Product")
    steps = [
        ("1", "Publish", "Services publish\nentity topics\nto Kafka", TEAL),
        ("2", "Materialize", "Data hydrated,\ndeduped, shaped\ninto clean entities", GREEN),
        ("3", "Fan Out", "Entities distributed\nto fit-for-purpose\nstorage lanes", ORANGE),
        ("4", "Compose", "GraphQL subgraphs\nfederate into\nsingle supergraph", ACCENT_PURPLE),
        ("5", "Serve", "Apollo Router\nserves unified\nexternal API", LIGHT_TEAL),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        left = Inches(0.4) + Inches(i * 2.55)
        add_shape(slide, left, Inches(1.7), Inches(2.3), Inches(2.0),
                  CARD_BG, border_color=color, border_width=2)
        add_text_box(slide, left, Inches(1.8), Inches(2.3), Inches(0.4),
                     num, font_size=30, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.15), Inches(2.3), Inches(0.3),
                     title, font_size=15, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.1), Inches(2.5), Inches(2.1), Inches(0.8),
                     desc, font_size=10, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        if i < 4:
            add_text_box(slide, left + Inches(2.3), Inches(2.2), Inches(0.25), Inches(0.4),
                         "→", font_size=22, font_color=MID_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    # Federation directives
    directives = [
        ("@key", "Entity Identity — primary key for federated entities"),
        ("@shareable", "Shared Fields — resolved by multiple subgraphs"),
        ("@external", "External Fields — owned by another subgraph"),
        ("extend type Query", "Root Query extension from subgraphs"),
    ]
    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(3), Inches(0.3),
                 "Key Federation Directives:", font_size=14, font_color=ORANGE, bold=True)
    for i, (d, desc) in enumerate(directives):
        y = Inches(4.4) + Inches(i * 0.35)
        add_text_box(slide, Inches(0.8), y, Inches(2.2), Inches(0.3),
                     d, font_size=12, font_color=TEAL, bold=True)
        add_text_box(slide, Inches(3.2), y, Inches(9), Inches(0.3),
                     desc, font_size=11, font_color=LIGHT_GRAY)

    add_text_box(slide, Inches(0.5), Inches(5.9), Inches(12), Inches(0.3),
                 "Gateway: ≤5ms p95  |  99.99% avail  |  100K RPS  |  Cache invalidation <1s  |  Rollback ≤5min",
                 font_size=11, font_color=ORANGE, alignment=PP_ALIGN.CENTER)

    # Subgraphs side by side
    make_section_slide(prs, "Subgraph Architecture\naperture-core & aperture-graph")

    slide = make_content_slide(prs, "Subgraph Architecture: Core & Graph")
    add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(4.5),
             "aperture-core  (Pointed Lookups)", [
                 "Backend: ScyllaDB (key-value lookups)",
                 "",
                 "Owns core entity types:",
                 "  • Identity, Entitlement, Source",
                 "  • AccessProfile, Role, Account",
                 "",
                 "All base properties:",
                 "  • name, email, department, etc.",
                 "  • Direct relationships: manager, accounts",
                 "",
                 "SLO: p95 ≤ 10ms  |  Freshness: < 60s",
                 "Scale: ≥ 100M identities, ≥ 5B records",
                 "Cost target: 30-50% reduction",
             ], icon="⚡", header_color=TEAL)

    add_card(slide, Inches(7), Inches(1.6), Inches(5.8), Inches(4.5),
             "aperture-graph  (Graph Traversals)", [
                 "Backend: JanusGraph on Scylla",
                 "",
                 "Owns graph-specific fields:",
                 "  • effectiveAccess, identitiesWithAccess",
                 "  • traverseIdentities, traverseEntitlements",
                 "  • searchIdentities (nested filtering)",
                 "",
                 "Relationship chain:",
                 "  identity → role → access_profile →",
                 "  entitlement → resource",
                 "",
                 "SLO: p95 < 1s (N-hop), < 100ms (10-hop)",
                 "Scale: ≥ 5B edges, ≥ 100M nodes",
             ], icon="🔗", header_color=GREEN)

    add_shape(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.45),
              CARD_BG, border_color=ORANGE, border_width=1)
    add_text_box(slide, Inches(0.5), Inches(6.32), Inches(12.3), Inches(0.4),
                 "🔄  Cross-Subgraph: Router stitches core + graph. One query → identity name (core) + effective access (graph)",
                 font_size=12, font_color=LIGHT_ORANGE, alignment=PP_ALIGN.CENTER)

    # Use Cases - 8 cards
    make_section_slide(prs, "Use Cases\nWhat Raptor Enables")

    slide = make_content_slide(prs, "Key Use Cases Powered by Raptor")
    use_cases = [
        ("🔍", "Identity Quick View", "Single query: name, email +\nall effective entitlements", TEAL),
        ("🔎", "Search with Context", "Search + profile info +\nreal-time activity counts", GREEN),
        ("🛡️", "Access Impact Analysis", "\"Who is affected if I\nrevoke this entitlement?\"", ORANGE),
        ("⚖️", "Policy Engine Context", "On-demand context for SoD,\nJIT, continuous compliance", ACCENT_PURPLE),
        ("🕸️", "Identity Graph", "Visual access exploration,\npath tracing, privilege analysis", LIGHT_TEAL),
        ("🔄", "Role Change Propagation", "Graph-scoped diffs →\ntargeted provision/deprovision", YELLOW),
        ("📊", "Real-Time Dashboards", "Live login counts, anomaly\ndetection without SF scans", ORANGE),
        ("🔐", "Simplified Materialization", "Pointed lookups replace\nstateful streaming joins", TEAL),
    ]
    for i, (icon, title, desc, color) in enumerate(use_cases):
        row, col = i // 4, i % 4
        left = Inches(0.4) + Inches(col * 3.15)
        top = Inches(1.6) + Inches(row * 2.5)
        add_card(slide, left, top, Inches(3.0), Inches(2.1),
                 title, desc.split("\n"), icon=icon, header_color=color)

    # Technology Decisions
    make_section_slide(prs, "Technology Decisions\nOptions & Trade-offs")

    slide = make_content_slide(prs, "Technology Selection Summary")
    decisions = [
        ("GraphQL Gateway", "Apollo Router + Federation 2", "✅ Excellent",
         "High-perf Rust gateway, mature federation,\npreview/compose, MCP integration", TEAL),
        ("Low-Latency Store", "ScyllaDB (self-hosted)", "✅ Excellent",
         "Ultra-low latency, cost-efficient,\nFedRAMP-ready, same cluster as graph", GREEN),
        ("Graph Database", "JanusGraph + Scylla", "✅ Excellent",
         "Mature open source, scalable,\nleverages same Scylla cluster", ORANGE),
        ("Real-Time Analytics", "ClickHouse", "✅ Excellent",
         "Extreme read speed, columnar storage,\nmature ecosystem, Kafka connectors", ACCENT_PURPLE),
        ("Alternative: Lookups", "Aerospike SaaS", "✅ Excellent",
         "Very low latency, built-in HA,\nbeing evaluated as alternative", LIGHT_TEAL),
    ]
    for i, (lane, choice, fit, notes, color) in enumerate(decisions):
        top = Inches(1.6) + Inches(i * 1.05)
        add_rect(slide, Inches(0.5), top, Inches(0.08), Inches(0.85), color)
        add_text_box(slide, Inches(0.8), top + Inches(0.05), Inches(2.5), Inches(0.3),
                     lane, font_size=14, font_color=WHITE, bold=True)
        add_text_box(slide, Inches(3.5), top + Inches(0.05), Inches(3.5), Inches(0.3),
                     choice, font_size=14, font_color=color, bold=True)
        add_text_box(slide, Inches(7.2), top + Inches(0.05), Inches(1.2), Inches(0.3),
                     fit, font_size=11, font_color=GREEN)
        add_text_box(slide, Inches(8.5), top + Inches(0.05), Inches(4.3), Inches(0.75),
                     notes, font_size=10, font_color=LIGHT_GRAY)

    make_section_slide(prs, "Q&A\n\n📧 dattu.marneni@sailpoint.com\n📋 DPINTAKE-33")

    prs.save("Raptor_Technical_Deep_Dive.pptx")
    print("✅  Created: Raptor_Technical_Deep_Dive.pptx")


# ==============================================================================
# PPT 3: ROADMAP & STATUS
# ==============================================================================
def create_roadmap_status():
    prs = new_prs()

    make_title_slide(prs, "🦅  Raptor\nRoadmap & Current Status")

    # Phases
    make_section_slide(prs, "Phased Roadmap\nPOC → Production")

    slide = make_content_slide(prs, "Four-Phase Delivery Roadmap")
    phases = [
        ("Phase I", "Prove It Works", "Q4'25–Q1'26", GREEN,
         ["Validate via POCs", "Observability dashboards", "Runbooks & cost models", "ALT decision gate"]),
        ("Phase II", "Golden Paths", "Q2'26", YELLOW,
         ["Kafka→Store→GraphQL docs", "Schema & SLO templates", "First product integration", "MySailPoint / Access Ctr"]),
        ("Phase III", "Self-Service", "Q3'26", TEAL,
         ["Prebuilt templates", "Onboarding guides", "Internal workshops", "Scale early adopters"]),
        ("Phase IV", "Scale & Harden", "Q4'26+", ACCENT_PURPLE,
         ["External GraphQL", "FedRAMP readiness", "Cost optimization", "SLO enforcement"]),
    ]
    for i, (phase, name, timeline, color, items) in enumerate(phases):
        left = Inches(0.4) + Inches(i * 3.15)
        add_shape(slide, left, Inches(1.6), Inches(3.0), Inches(4.3),
                  CARD_BG, border_color=color, border_width=2)
        add_rect(slide, left, Inches(1.6), Inches(3.0), Inches(0.06), color)
        add_text_box(slide, left, Inches(1.75), Inches(3.0), Inches(0.3),
                     phase, font_size=13, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.05), Inches(3.0), Inches(0.35),
                     name, font_size=17, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.4), Inches(3.0), Inches(0.25),
                     timeline, font_size=11, font_color=color, alignment=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            add_text_box(slide, left + Inches(0.15), Inches(2.8) + Inches(j * 0.35),
                         Inches(2.7), Inches(0.3),
                         f"• {item}", font_size=11, font_color=LIGHT_GRAY)

    add_text_box(slide, Inches(0.4), Inches(6.1), Inches(3.0), Inches(0.3),
                 "◀ WE ARE HERE", font_size=13, font_color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    # DACI Status
    make_section_slide(prs, "DACI Status\nDecision Progress Across Lanes")

    slide = make_content_slide(prs, "DACI Decision Status & POC Progress")
    dacis = [
        ("🌐", "GraphQL Layer", "Apollo Router + Federation 2", "📝 Draft", "SAASD8NG-4502", "Router p95 ≤ 5ms;\nfirst subgraph integrated", TEAL),
        ("⚡", "Pointed Lookups", "ScyllaDB (self-hosted + Cloud)", "🔄 In Progress", "SAASD8NG-4501", "p95 ≤ 10ms @ 5K RPS;\ncost ↓ 30%", GREEN),
        ("🔗", "Graph Database", "JanusGraph on Scylla", "🔄 In Progress", "SAASD8NG-4273", "p95 < 1s N-hop;\n5B edges scale", ORANGE),
        ("📊", "Real-Time Analytics", "ClickHouse (SH/SaaS)", "📝 Draft", "SAASD8NG-4972", "p95 < 300ms;\ncost ↓ 40% vs SF", ACCENT_PURPLE),
        ("🔍", "Search Layer", "Aperture Search subgraph", "📝 Draft", "—", "Search 2.0 strategy;\ncontract parity", LIGHT_TEAL),
    ]
    for i, (icon, name, tech, status, epic, criteria, color) in enumerate(dacis):
        top = Inches(1.6) + Inches(i * 1.0)
        add_rect(slide, Inches(0.5), top, Inches(0.08), Inches(0.85), color)
        add_text_box(slide, Inches(0.7), top + Inches(0.05), Inches(0.4), Inches(0.3),
                     icon, font_size=16, font_color=color, bold=True)
        add_text_box(slide, Inches(1.2), top + Inches(0.05), Inches(2.0), Inches(0.3),
                     name, font_size=14, font_color=WHITE, bold=True)
        add_text_box(slide, Inches(3.4), top + Inches(0.05), Inches(2.5), Inches(0.3),
                     tech, font_size=11, font_color=LIGHT_GRAY)
        sc = GREEN if "Progress" in status else YELLOW
        add_text_box(slide, Inches(6.1), top + Inches(0.05), Inches(1.5), Inches(0.3),
                     status, font_size=11, font_color=sc, bold=True)
        add_text_box(slide, Inches(7.8), top + Inches(0.05), Inches(1.8), Inches(0.3),
                     epic, font_size=10, font_color=TEAL)
        add_text_box(slide, Inches(9.8), top + Inches(0.05), Inches(3.2), Inches(0.75),
                     criteria, font_size=10, font_color=LIGHT_GRAY)

    add_text_box(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.3),
                 "⚠️  Driver: Derrick Mink (departed) → Ownership gap  |  Approvers: Dan Sparks, Fuad Rashid  |  INIT-2073: In Progress",
                 font_size=10, font_color=ORANGE, alignment=PP_ALIGN.CENTER)

    # Risks & Actions
    make_section_slide(prs, "Risks & Actions\nMitigations & Next Steps")

    slide = make_content_slide(prs, "Risks, Mitigations & Immediate Actions")
    risks = [
        ("⚠️", "Driver Departed", "Derrick Mink left — gap\nacross all DACI lanes", "Assign new drivers\nper lane immediately", RED),
        ("🚧", "DevOps Bottleneck", "Infrastructure intake\ncan block POC velocity", "Pre-stub stories,\nalign early with SRE", ORANGE),
        ("🏗️", "Infra Complexity", "Self-hosted DBs need\nops expertise", "Evaluate SaaS\nalternatives in parallel", YELLOW),
        ("👥", "Team Capacity", "Limited bandwidth for\nmulti-lane POCs", "Engage evolv/tiger\nteams for env setup", TEAL),
    ]
    for i, (icon, title, risk, mitigation, color) in enumerate(risks):
        left = Inches(0.4) + Inches(i * 3.15)
        add_card(slide, left, Inches(1.6), Inches(3.0), Inches(2.3),
                 title, (risk + "\n\n✅ " + mitigation).split("\n"),
                 icon=icon, header_color=color)

    # Next steps
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(5), Inches(0.3),
                 "⚡ Immediate Next Steps:", font_size=16, font_color=ORANGE, bold=True)
    actions = [
        ("1️⃣", "Assess current POC state — what's deployed & running"),
        ("2️⃣", "Clarify ownership — new drivers for each DACI lane"),
        ("3️⃣", "Address DPINTAKE-33 — delivery date for Identity Graph"),
        ("4️⃣", "Unblock backlog — assign subgraph stories 5980/5981/5982"),
        ("5️⃣", "Rename to Raptor — update Confluence, Jira, Slack"),
        ("6️⃣", "Establish weekly cadence — \"Raptor Updates\" posts"),
        ("7️⃣", "Schedule kickoff sync — Forge, Keel, DIPO, Double Zero"),
    ]
    for i, (num, action) in enumerate(actions):
        col = i // 4
        row = i % 4
        left = Inches(0.5) + Inches(col * 6.3)
        top = Inches(4.6) + Inches(row * 0.45)
        add_text_box(slide, left, top, Inches(6.0), Inches(0.4),
                     f"{num}  {action}", font_size=12, font_color=LIGHT_GRAY)

    make_section_slide(prs, "Q&A\n\n📧 dattu.marneni@sailpoint.com\n📋 DPINTAKE-33\n💬 #raptor-data-platform")

    prs.save("Raptor_Roadmap_and_Status.pptx")
    print("✅  Created: Raptor_Roadmap_and_Status.pptx")


if __name__ == "__main__":
    create_exec_overview()
    create_tech_deep_dive()
    create_roadmap_status()
    print("\n🎉  All 3 presentations created — SailPoint template + rich visual content!")
