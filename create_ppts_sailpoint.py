#!/usr/bin/env python3
"""Generate Raptor presentations using the SailPoint branded template."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

TEMPLATE = "/Users/dattu.marneni/Desktop/Sun Aug 31 16-53-08 2025/Atlan_Beta_User_Training.pptx"

BLUE_ACCENT = RGBColor(0x20, 0x26, 0xD2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x66, 0x66, 0x66)
TEAL = RGBColor(0x00, 0xB4, 0xD8)
GREEN = RGBColor(0x28, 0xA7, 0x45)
ORANGE = RGBColor(0xF7, 0x7F, 0x00)
RED = RGBColor(0xDC, 0x35, 0x45)
PURPLE = RGBColor(0x6F, 0x42, 0xC1)

# Layout indices from template
LAYOUT_TITLE = 0        # Title Slide-Core Patterns
LAYOUT_SECTION = 20     # Section Header
LAYOUT_SECTION_DARK = 23  # Section Header Dark
LAYOUT_CONTENT = 7      # Title and Content
LAYOUT_CONTENT_DARK = 8 # Title and Content - Dark gradient
LAYOUT_BLANK = 25       # Blank
LAYOUT_AGENDA = 5       # Agenda


def new_prs():
    """Create a new presentation from the SailPoint template with all content slides removed."""
    prs = Presentation(TEMPLATE)
    # Remove existing slides by manipulating the XML directly
    from lxml import etree
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId is None:
            rId = sldId.get('r:id')
        if rId:
            try:
                prs.part.drop_rel(rId)
            except (KeyError, Exception):
                pass
        sldIdLst.remove(sldId)
    return prs


def add_title_slide(prs, title_text, subtitle_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:  # Title
            shape.text = ""
            tf = shape.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title_text
            run.font.name = "Poppins SemiBold"
            run.font.size = Pt(40)
            run.font.color.rgb = WHITE
    return slide


def add_section_slide(prs, title_text, speaker=""):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION_DARK])
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:  # Title
            shape.text = ""
            tf = shape.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title_text
            run.font.name = "Poppins SemiBold"
            run.font.size = Pt(36)
            run.font.color.rgb = WHITE
        if speaker and hasattr(shape.placeholder_format, 'idx') and shape.placeholder_format.idx == 1:
            shape.text = f"Speaker:  {speaker}"
    return slide


def add_content_slide(prs, title_text, bullets=None, sub_bullets=None, dark=False):
    layout_idx = LAYOUT_CONTENT_DARK if dark else LAYOUT_CONTENT
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:  # Title
            shape.text = ""
            tf = shape.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title_text
            run.font.name = "Poppins SemiBold"
            run.font.size = Pt(28)
            run.font.color.rgb = WHITE if dark else DARK_TEXT

        if shape.placeholder_format.idx in (1, 10) and bullets:  # Content
            shape.text = ""
            tf = shape.text_frame
            text_color = WHITE if dark else DARK_TEXT
            heading_color = WHITE if dark else BLACK

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                is_heading = bullet.startswith("##")
                is_sub = bullet.startswith("  ")
                is_quote = bullet.startswith('"')
                is_blank = bullet.strip() == ""

                if is_blank:
                    p.text = ""
                    p.space_before = Pt(6)
                    p.space_after = Pt(6)
                    continue

                clean = bullet.lstrip("# ").lstrip("  ")

                run = p.add_run()
                run.text = clean
                run.font.name = "Poppins"
                run.font.size = Pt(14 if is_sub else (20 if is_heading else 16))
                run.font.bold = is_heading or (not is_sub and not is_quote)
                run.font.color.rgb = BLUE_ACCENT if is_heading else (LIGHT_GRAY if is_sub else (text_color if is_quote else heading_color))

                p.space_before = Pt(12 if is_heading else 4)
                p.space_after = Pt(4)
                p.level = 1 if is_sub else 0

                if sub_bullets and bullet in sub_bullets:
                    for sb in sub_bullets[bullet]:
                        sp = tf.add_paragraph()
                        sr = sp.add_run()
                        sr.text = sb
                        sr.font.name = "Poppins"
                        sr.font.size = Pt(14)
                        sr.font.color.rgb = text_color
                        sr.font.bold = False
                        sp.level = 1
                        sp.space_before = Pt(2)
                        sp.space_after = Pt(2)

    # Update footer
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 12:  # Footer
            shape.text = "© 2026 SailPoint Technologies, Inc. All rights reserved."

    return slide


def add_agenda_slide(prs, items):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = ""
            tf = shape.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Agenda"
            run.font.name = "Poppins SemiBold"
            run.font.size = Pt(28)
            run.font.color.rgb = DARK_TEXT

        if shape.placeholder_format.idx in (1, 10):
            shape.text = ""
            tf = shape.text_frame
            for i, (num, title, desc) in enumerate(items):
                if i > 0:
                    p_blank = tf.add_paragraph()
                    p_blank.text = ""
                    p_blank.space_before = Pt(8)

                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                run = p.add_run()
                run.text = f"{num}    {title}"
                run.font.name = "Poppins"
                run.font.size = Pt(20)
                run.font.bold = True
                run.font.color.rgb = BLACK

                p2 = tf.add_paragraph()
                run2 = p2.add_run()
                run2.text = f"       {desc}"
                run2.font.name = "Poppins"
                run2.font.size = Pt(14)
                run2.font.color.rgb = LIGHT_GRAY
                p2.level = 0
                p2.space_before = Pt(2)

    return slide


# ==============================================================================
# PPT 1: EXECUTIVE OVERVIEW
# ==============================================================================
def create_exec_overview():
    prs = new_prs()

    # 1. Title
    add_title_slide(prs, "Raptor\nOne Door to Data at Scale")

    # 2. Agenda
    add_agenda_slide(prs, [
        ("01", "The Challenge", "Why change is needed now"),
        ("02", "What is Raptor?", "Vision and approach"),
        ("03", "Architecture", "Five specialized data lanes"),
        ("04", "Business Impact", "Target outcomes and SLOs"),
        ("05", "Next Steps", "Immediate actions and timeline"),
    ])

    # 3. Section: The Challenge
    add_section_slide(prs, "The Challenge\nWhy We Need Raptor")

    # 4. Problem - Cost
    add_content_slide(prs, "The Challenge: Unsustainable Architecture", [
        "## Current Pain Points",
        "",
        "💰  Unsustainable cost curve",
        "  Search materialization costs exceed $150,000/month in compute alone",
        "  Re-materialization multiplies cost linearly with data volume",
        "",
        "🕸️  Growing complexity",
        "  Each product feature adds bespoke pipelines, schema joins, and edge-case handling",
        "  What began as flexible has evolved into a brittle web of dependencies",
        "",
        "🚧  Platform bottleneck",
        "  Data platform is the critical path for too many teams",
        "  Even small product changes require deep data engineering involvement",
        "",
        "⏰  Fragile freshness",
        "  \"Real-time\" today means frequent batch runs and risk-prone replays",
        "  Any schema tweak ripples through a brittle chain of dependencies",
    ], dark=True)

    # 5. Section: What is Raptor
    add_section_slide(prs, "What is Raptor?")

    # 6. Raptor Vision
    add_content_slide(prs, "Raptor: A Single Door to Data at Scale", [
        "## A single external API that powers every product experience needing data",
        "",
        "Instead of pre-building massive documents through expensive materializers,",
        "  data is stored once and composed dynamically at query time",
        "  through a unified GraphQL API that routes to the right backend",
        "",
        "## Before → After",
        "",
        "❌  Big pre-built documents per use case  →  ✅  On-demand composition via federated subgraphs",
        "❌  Expensive backfills & replays  →  ✅  Real-time by design: fetch fresh state directly",
        "❌  Multi-day lead time for new fields  →  ✅  New fields show up in hours, not weeks",
        "❌  Freshness via constant re-materialization  →  ✅  Materializers become thin or optional",
        "",
        "\"Think of it like a restaurant: today we cook every item on the menu just in case.",
        "  With Raptor, we prepare ingredients once and cook dishes fresh when ordered.\"",
    ])

    # 7. Core Principles
    add_content_slide(prs, "Core Principles", [
        "🎯  Right data, right store",
        "  Each access pattern belongs in the backend designed for it — no more bending one system to do everything",
        "",
        "🚪  One API, many backends",
        "  A single external API presents data consistently, while teams own their federated subgraphs",
        "",
        "🌊  Streaming first",
        "  Data is cleaned, deduped, and hydrated once, then fanned out everywhere it's needed",
        "",
        "🏢  Tenant-first design",
        "  Every layer enforces isolation, limits, and observability per org or pod",
        "",
        "🚀  Progress over perfection",
        "  Start with POCs and vertical slices, then harden iteratively as adoption grows",
        "",
        "📡  Serve in place when possible",
        "  If a domain API already runs on a scalable backend, expose a GraphQL subgraph directly",
    ])

    # 8. Section: Architecture
    add_section_slide(prs, "Architecture\nFive Specialized Data Lanes")

    # 9. Architecture Lanes
    add_content_slide(prs, "Architecture: Five Specialized Data Lanes + Unified API", [
        "## ⚡  Pointed Lookups  —  ScyllaDB / Aerospike",
        "  Entity hydration in single-digit milliseconds  |  p95 ≤ 10ms  |  Scale: ≥ 5B records",
        "",
        "## 🔗  Graph Database  —  JanusGraph on Scylla",
        "  Relationship traversals: identity → role → access profile → entitlement  |  p95 < 1s",
        "",
        "## 📊  Real-Time Aggregates  —  ClickHouse",
        "  Dashboards, anomaly detection, usage metrics  |  p95 < 300ms  |  > 5B events/day",
        "",
        "## 🔍  Search & Discovery  —  OpenSearch / Elasticsearch",
        "  Full-text search, autocomplete, purpose-built indexes  |  p95 < 200ms",
        "",
        "## 📚  Historical & Batch  —  Snowflake",
        "  Large-scale reporting, compliance, ML training  |  Hours-level freshness",
        "",
        "🌐  All lanes unified behind Apollo Router + GraphQL Federation 2  →  Single External Schema",
    ], dark=True)

    # 10. Section: Business Impact
    add_section_slide(prs, "Business Impact\nTarget Outcomes")

    # 11. Outcomes
    add_content_slide(prs, "Target Outcomes & Business Impact", [
        "## Key Metrics",
        "",
        "💰  30–50% cost reduction vs. Snowflake baseline",
        "",
        "🚀  < 2 days to surface new entity via GraphQL (currently weeks)",
        "",
        "🗑️  ≥ 3 materializers retired by mid-2026",
        "",
        "📊  99.99% availability target",
        "",
        "😊  ≥ 8/10 developer satisfaction score",
        "",
        "✅  ≥ 75% POC → Golden Path conversion rate",
        "",
        "## What This Enables for Product Teams",
        "",
        "  Faster delivery: no custom pipelines or multi-week engineering work",
        "  Real-time data: fresh identity details, updated entitlements, live usage counts",
        "  Lower cost: reduce data duplication and rebuild cycles",
        "  AI-ready: GraphQL + MCP enables AI copilots and natural language querying",
    ])

    # 12. Section: Next Steps
    add_section_slide(prs, "Next Steps\nImmediate Actions")

    # 13. Actions
    add_content_slide(prs, "Immediate Next Steps — Raptor Transition", [
        "1️⃣  Assess current POC state",
        "  Review what is deployed and running in dev (Apollo Server, ScyllaDB, JanusGraph)",
        "",
        "2️⃣  Clarify ownership",
        "  With original driver's departure, establish new drivers for each DACI lane",
        "",
        "3️⃣  Address DPINTAKE-33",
        "  Provide delivery date and incremental plan for Identity Graph use case",
        "",
        "4️⃣  Unblock the backlog",
        "  Assign and prioritize subgraph stories (SAASD8NG-5980/5981/5982)",
        "",
        "5️⃣  Rename to Raptor",
        "  Update Confluence pages, Jira labels/epics, and Slack channels",
        "",
        "6️⃣  Establish weekly cadence",
        "  Start \"Raptor Updates\" posts in Slack and Confluence",
        "",
        "7️⃣  Schedule kickoff sync with Forge, Keel, DIPO, and Double Zero",
    ], dark=True)

    # 14. Q&A
    add_section_slide(prs, "Q&A")

    prs.save("Raptor_Executive_Overview.pptx")
    print("✅  Created: Raptor_Executive_Overview.pptx")


# ==============================================================================
# PPT 2: TECHNICAL DEEP DIVE
# ==============================================================================
def create_tech_deep_dive():
    prs = new_prs()

    add_title_slide(prs, "Raptor\nTechnical Deep Dive")

    add_agenda_slide(prs, [
        ("01", "GraphQL Federation", "How the unified API works"),
        ("02", "Subgraph Architecture", "aperture-core & aperture-graph"),
        ("03", "Use Cases", "What Raptor enables"),
        ("04", "Technology Decisions", "Options and trade-offs"),
        ("05", "Infrastructure", "Dev environment & repos"),
    ])

    # GraphQL Federation
    add_section_slide(prs, "GraphQL Federation\nHow the Unified API Works")

    add_content_slide(prs, "Data Flow: Source → Lane → GraphQL → Product", [
        "## Step 1: Publish",
        "  Authoritative services publish into entity topics (Kafka)",
        "",
        "## Step 2: Materialize Once",
        "  Data hydrated, deduped, and shaped into clean entity documents",
        "",
        "## Step 3: Fan Out",
        "  Entities distributed to fit-for-purpose storage lanes",
        "",
        "## Step 4: Compose",
        "  GraphQL subgraphs federate into single supergraph via Apollo Router",
        "",
        "## Step 5: Serve",
        "  Single API serves all product experiences — one query, multiple backends",
        "",
        "Technical Requirements:",
        "  Gateway overhead: ≤ 5ms p95  |  Availability: ≥ 99.99%  |  Scale: 100K RPS",
        "  Cache invalidation: < 1s  |  Schema rollback: ≤ 5min  |  Standard SailPoint auth",
    ], dark=True)

    add_content_slide(prs, "Federation Key Concepts", [
        "## @key  —  Entity Identity",
        "  Defines primary key for federated entities: type Identity @key(fields: \"id\")",
        "",
        "## @shareable  —  Shared Fields",
        "  Field can be resolved by multiple subgraphs (e.g., id, name)",
        "",
        "## @external  —  External Fields",
        "  Declares field is owned by another subgraph — router handles stitching",
        "",
        "## extend type Query  —  Root Query Extension",
        "  Each subgraph adds its own query fields to the unified schema",
        "",
        "## Entity Stubs vs Full Definitions",
        "  aperture-core: full definition of Identity (name, email, department, etc.)",
        "  aperture-graph: entity stub — adds only effectiveAccess, graph traversals",
        "  Router merges both so clients query all fields in one request",
        "",
        "## Query Planning",
        "  Router creates execution plan → fetches from each subgraph → merges → returns",
    ])

    # Subgraphs
    add_section_slide(prs, "Subgraph Architecture\naperture-core & aperture-graph")

    add_content_slide(prs, "aperture-core  —  Pointed Lookups (ScyllaDB)", [
        "## What It Owns",
        "  Core entity types: Identity, Entitlement, Source, AccessProfile, Role, Account",
        "  All base properties: name, email, department, isManager, etc.",
        "  Direct relationships: entitlementAssignments, manager, accounts",
        "",
        "## How It Works",
        "  Entities retrieved by composite keys in single-digit milliseconds",
        "  Option to defer loading large document fragments",
        "  Event-driven partial updates using JSON patch semantics",
        "",
        "## Technical Targets",
        "  Latency: p95 ≤ 10ms  |  Freshness: < 60s from Kafka event",
        "  Availability: ≥ 99.99%  |  Scale: linear, per-tenant partitioning",
        "  Cost: 30–50% lower than Elasticsearch/Snowflake hydration",
        "",
        "## Key Benefit",
        "  Identity document becomes a pure function of current state — not a product of historical join order",
    ], dark=True)

    add_content_slide(prs, "aperture-graph  —  Graph Traversals (JanusGraph)", [
        "## What It Owns",
        "  Graph-specific fields: effectiveAccess, identitiesWithAccess",
        "  Graph traversal queries: traverseIdentities, traverseEntitlements",
        "  Graph-based search: searchIdentities (with nested filtering + AMM metadata)",
        "",
        "## Relationship Chain",
        "  identity → role → access_profile → entitlement → resource",
        "",
        "## How It Works",
        "  On-demand traversal over live relationships with predictable latency",
        "  Stream-driven edge creation from Kafka entity topics",
        "  Schema-driven ingestion and dynamic edge creation",
        "",
        "## Technical Targets",
        "  Traversal: p95 ≤ 1s (N-hop)  |  Targeted lookups: p95 ≤ 100ms (10-hop)",
        "  Freshness: < 60s  |  Scale: ≥ 5B edges across ≥ 100M nodes",
        "",
        "## Key Benefit",
        "  Enables real-time role change propagation, policy evaluation, and JIT access decisions",
    ])

    # Use Cases
    add_section_slide(prs, "Use Cases\nWhat Raptor Enables")

    add_content_slide(prs, "Use Cases: Search, Identity & Policy", [
        "## 🔍  Simplified Identity Materialization",
        "  Replace stateful streaming joins with stateless pointed lookups",
        "  Identity docs become enrichment operations, not join products",
        "  Complexity stays flat as new relationships are added",
        "",
        "## ⚖️  Identity Context for Policy Engine",
        "  On-demand identity context for SoD, JIT access, continuous compliance",
        "  No pre-materialized identity document required",
        "  Supports on-demand, reactive, and scheduled evaluation modes",
        "",
        "## 🕸️  Knowledge & Identity Graph",
        "  Visual exploration of access relationships — nodes and edges",
        "  Interactive path tracing: \"how does this identity have this access?\"",
        "  Privilege exposure analysis and impact forecasting",
        "  Graph views computed at query time — no pre-materialized graph datasets",
    ], dark=True)

    add_content_slide(prs, "Use Cases: Product Experiences", [
        "## ⚡  Identity Quick View",
        "  One query: name, email, org + all effective entitlements — no separate lookups",
        "",
        "## 🔎  Search with Context",
        "  Autocomplete + search results enriched with profile info and live activity counts",
        "",
        "## 🛡️  Access Impact Analysis",
        "  \"If I revoke this entitlement, who is affected and how have they used it?\"",
        "",
        "## 📊  Real-Time Dashboards",
        "  Live login counts, access requests, anomaly detection without Snowflake scans",
        "",
        "## 🔄  Role Change Propagation",
        "  Graph-scoped diffs for role changes → targeted provision/deprovision",
        "",
        "## 🤖  AI-Ready via MCP",
        "  Schema auto-published for AI agents — natural language querying",
    ])

    # Technology
    add_section_slide(prs, "Technology Decisions\nOptions & Trade-offs")

    add_content_slide(prs, "Technology Selection Summary", [
        "## 🌐  GraphQL Gateway  →  Apollo Router + Federation 2  ✅ Excellent",
        "  High-perf Rust gateway, mature federation, preview/compose, MCP integration",
        "",
        "## ⚡  Low-Latency Store  →  ScyllaDB (self-hosted)  ✅ Excellent",
        "  Ultra-low latency, cost-efficient, FedRAMP-ready",
        "  Alternative: Aerospike SaaS — very low latency, being evaluated (SAASD8NG-5554)",
        "",
        "## 🔗  Graph Database  →  JanusGraph + Scylla  ✅ Excellent",
        "  Mature open source, scalable, leverages same Scylla cluster as lookup lane",
        "",
        "## 📊  Real-Time Analytics  →  ClickHouse  ✅ Excellent",
        "  Extreme read speed, columnar storage, mature Kafka connectors",
        "",
        "## 🔍  Search Layer  →  Decision Pending",
        "  Aperture Search subgraph vs optimized materializer model",
        "  Must preserve existing customer search contract (filters, freshness, paging)",
    ], dark=True)

    # Infrastructure
    add_section_slide(prs, "Infrastructure\nDev Environment & Repos")

    add_content_slide(prs, "GitHub Repositories & Dev Environment", [
        "## GitHub Repositories",
        "  sailpoint-core/aperture-core — ScyllaDB pointed lookups service",
        "  sailpoint-core/aperture-graph — JanusGraph graph traversals service",
        "  sailpoint-core/aperture-search — Elasticsearch search service",
        "  sailpoint-core/pointed-lookups-dip — DIP for pointed lookups pipeline",
        "  sailpoint/gitops-k8s (apollo-poc branch) — Apollo Server deployment",
        "",
        "## Dev Environment",
        "  Apollo Server: deployed via ArgoCD in dev-us-east-1",
        "  ScyllaDB: 3-node cluster in AWS us-east-1 (Cloud)",
        "  Mach5 UI: https://mach5-m5s.odin-use1-2.ida.cloud.sailpoint.com/",
        "  Grafana: Mach5 Health Metrics + Entity Metrics dashboards",
        "",
        "## Jira Epics",
        "  SAASD8NG-4502: Apollo Server GraphQL",
        "  SAASD8NG-4501: Pointed Lookup Pipeline (ScyllaDB)",
        "  SAASD8NG-4273: Graph Database POC (JanusGraph)",
        "  SAASD8NG-4972: Real-Time Analytics POC (ClickHouse)",
    ])

    add_section_slide(prs, "Q&A")

    prs.save("Raptor_Technical_Deep_Dive.pptx")
    print("✅  Created: Raptor_Technical_Deep_Dive.pptx")


# ==============================================================================
# PPT 3: ROADMAP & STATUS
# ==============================================================================
def create_roadmap_status():
    prs = new_prs()

    add_title_slide(prs, "Raptor\nRoadmap & Current Status")

    add_agenda_slide(prs, [
        ("01", "Phased Roadmap", "POC → Golden Path → Self-Service → Scale"),
        ("02", "DACI Status", "Decision status across all lanes"),
        ("03", "Current Progress", "What's done, what's in flight"),
        ("04", "Risks & Actions", "Mitigations and immediate next steps"),
    ])

    # Phase I
    add_section_slide(prs, "Phased Roadmap\nPOC → Production")

    add_content_slide(prs, "Four-Phase Delivery Roadmap", [
        "## 🟢  Phase I — Prove It Works  (Q4'25 – Q1'26)  ← CURRENT",
        "  Validate architecture through fast, focused POCs across all four DACIs",
        "  Joint observability dashboards | Runbooks & cost models | ALT decision gate",
        "",
        "## 🟡  Phase II — Golden Paths  (Q2'26)",
        "  Convert POC results into production-ready reference implementations",
        "  Documented Kafka → Store → GraphQL flows | Schema & SLO templates",
        "  First product integration (MySailPoint or Access Center)",
        "",
        "## 🔵  Phase III — Self-Service & Adoption  (Q3'26)",
        "  Enable broader adoption through prebuilt templates and onboarding guides",
        "  Internal workshops | Cost dashboards | Scale early adopters",
        "",
        "## 🟣  Phase IV — Scale & Hardening  (Q4'26+)",
        "  External GraphQL exposure | FedRAMP readiness",
        "  Cost optimization | SLO enforcement | Full self-service platform",
    ], dark=True)

    # DACI Status
    add_section_slide(prs, "DACI Status\nDecision Progress Across Lanes")

    add_content_slide(prs, "DACI Decision Status — All Lanes", [
        "## 🌐  GraphQL Layer  —  📝 Draft",
        "  Apollo Router + Federation 2 | SAASD8NG-4502",
        "  Success: Router p95 overhead ≤ 5ms; first subgraph integrated",
        "",
        "## ⚡  Pointed Lookups  —  🔄 In Progress",
        "  ScyllaDB (self-hosted + Cloud) | SAASD8NG-4501",
        "  Success: p95 ≤ 10ms at 5K RPS/pod; cost ↓ 30%",
        "",
        "## 🔗  Graph Database  —  🔄 In Progress",
        "  JanusGraph on Scylla | SAASD8NG-4273",
        "  Success: p95 < 1s for N-hop traversals; linear scale at 5B edges",
        "",
        "## 📊  Real-Time Analytics  —  📝 Draft",
        "  ClickHouse (self-hosted / SaaS) | SAASD8NG-4972",
        "  Success: p95 < 300ms; cost ↓ 40% vs Snowflake",
        "",
        "## 🔍  Search Layer  —  📝 Draft",
        "  Aperture Search subgraph | Decision: proceed, pivot, or defer",
        "  Must preserve existing customer search contract (REST → GraphQL adapter)",
    ])

    # Progress
    add_section_slide(prs, "Current Progress\nWhat's Done, What's In Flight")

    add_content_slide(prs, "Phase I Action Tracking", [
        "## ✅  Completed",
        "  Separate DACI created for each workstream (GraphQL, Lookups, Aggregates, Graph)",
        "  POC epics created for Phase I (one per lane)",
        "  Shared Slack channel created for async collaboration",
        "",
        "## 🟡  In Progress",
        "  INIT-2073: Ensuring INIT created and communicated to business",
        "  Identifying POC resources (contractors or Forge engineers)",
        "",
        "## ⬜  Not Started",
        "  POC Exit Review template",
        "  Centralized POC tracking view",
        "  Aperture POC kickoff sync with Forge, Keel, DIPO, Double Zero",
        "  All-hands presentation introducing Raptor vision",
        "  Weekly \"Raptor Updates\" posts",
        "",
        "## ⚠️  Key Gap",
        "  Original driver (Derrick Mink) has departed — ownership to be reassigned",
    ], dark=True)

    add_content_slide(prs, "Key Jira Tickets — Status", [
        "## 🔥  Active",
        "  DPINTAKE-33: Project Aperture Delivery for Identity Graph (Assigned: Dattu Marneni)",
        "  INIT-2073: Aperture INIT — business communication (In Progress)",
        "",
        "## 📋  Backlog — Subgraph Creation",
        "  SAASD8NG-5982: Create subgraph for Aperture-Search (ElasticSearch) — Unassigned",
        "  SAASD8NG-5981: Create subgraph for Aperture-Graph (JanusGraph) — Unassigned",
        "  SAASD8NG-5980: Create subgraph for Aperture-Core (ScyllaDB) — Unassigned",
        "",
        "## 📋  Backlog — Infrastructure",
        "  DPE-179: Service User Creation for aperture-core — Ashish Pandita",
        "  DPE-165: Service User Creation for aperture-graph — Subham Jain",
        "  DPE-128: Deprecate deploy_svc_user from aperture-graph — Subham Jain",
        "",
        "## 📋  Backlog — Alternatives",
        "  SAASD8NG-5554: Pointed Lookups on Aerospike — Chris Schneider",
    ])

    # Risks
    add_section_slide(prs, "Risks & Actions\nMitigations & Next Steps")

    add_content_slide(prs, "Risks & Mitigations", [
        "## ⚠️  Driver Departed",
        "  Risk: Derrick Mink left — ownership gap across all DACI lanes",
        "  Mitigation: Assign new drivers per lane immediately; clarify with leadership",
        "",
        "## 🚧  DevOps Intake Bottlenecks",
        "  Risk: Infrastructure intake can block POC velocity",
        "  Mitigation: Pre-stub stories, align early with SRE teams",
        "",
        "## 🏗️  Infrastructure Complexity",
        "  Risk: Self-hosted databases require significant ops expertise",
        "  Mitigation: Evaluate SaaS alternatives in parallel (Scylla Cloud, Aerospike)",
        "",
        "## 👥  Team Capacity Constraints",
        "  Risk: Limited bandwidth for multi-lane POCs simultaneously",
        "  Mitigation: Engage evolv Consulting / tiger teams for environment setup",
        "",
        "## 📈  Adoption Lag",
        "  Risk: Enterprise drag — gradual adoption may be slow",
        "  Mitigation: Start with high-visibility golden paths and demos",
    ], dark=True)

    add_content_slide(prs, "Immediate Next Steps — Raptor Transition", [
        "1️⃣  Assess current POC state",
        "  Review deployed environments: Apollo Server, ScyllaDB, JanusGraph",
        "",
        "2️⃣  Clarify ownership per DACI lane",
        "  Establish new drivers; update DACI documents with Raptor branding",
        "",
        "3️⃣  Address DPINTAKE-33",
        "  Reasonable delivery date + incremental plan for Identity Graph",
        "",
        "4️⃣  Unblock the backlog",
        "  Assign subgraph stories (5980/5981/5982) and service user creation",
        "",
        "5️⃣  Rename to Raptor across all systems",
        "  Confluence, Jira labels/epics, Slack channels",
        "",
        "6️⃣  Establish weekly cadence — \"Raptor Updates\"",
        "",
        "7️⃣  Schedule kickoff sync with Forge, Keel, DIPO, Double Zero leads",
    ])

    add_section_slide(prs, "Q&A")

    prs.save("Raptor_Roadmap_and_Status.pptx")
    print("✅  Created: Raptor_Roadmap_and_Status.pptx")


if __name__ == "__main__":
    create_exec_overview()
    create_tech_deep_dive()
    create_roadmap_status()
    print("\n🎉  All 3 presentations created in SailPoint template style!")
