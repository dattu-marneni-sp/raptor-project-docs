#!/usr/bin/env python3
"""Generate SailPoint-styled PowerPoint presentations for the Raptor project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# SailPoint brand colors
NAVY = RGBColor(0x0D, 0x1B, 0x3E)
DARK_BLUE = RGBColor(0x14, 0x27, 0x52)
SAILPOINT_BLUE = RGBColor(0x00, 0x33, 0x66)
TEAL = RGBColor(0x00, 0xB4, 0xD8)
LIGHT_TEAL = RGBColor(0x48, 0xCA, 0xE4)
ORANGE = RGBColor(0xF7, 0x7F, 0x00)
LIGHT_ORANGE = RGBColor(0xFF, 0xA6, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
MID_GRAY = RGBColor(0x8C, 0x8C, 0xA0)
DARK_GRAY = RGBColor(0x3A, 0x3A, 0x4A)
GREEN = RGBColor(0x28, 0xA7, 0x45)
YELLOW = RGBColor(0xFF, 0xC1, 0x07)
RED = RGBColor(0xDC, 0x35, 0x45)
ACCENT_PURPLE = RGBColor(0x6F, 0x42, 0xC1)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=14, font_color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4), font_name="Calibri"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet_text(text_frame, text, font_size=14, font_color=WHITE, bold=False, level=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.level = level
    p.space_before = Pt(4)
    p.space_after = Pt(4)
    return p


def create_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)

    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), TEAL)

    # Bottom accent bar
    add_rect(slide, Inches(0), Inches(7.0), SLIDE_WIDTH, Inches(0.08), ORANGE)

    # Left accent stripe
    add_rect(slide, Inches(0), Inches(0.08), Inches(0.15), Inches(6.92), TEAL)

    # Title
    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                 title, font_size=44, font_color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

    # Subtitle
    add_text_box(slide, Inches(1.5), Inches(3.8), Inches(9), Inches(1.0),
                 subtitle, font_size=22, font_color=LIGHT_TEAL, bold=False, alignment=PP_ALIGN.LEFT)

    # Logo area
    add_text_box(slide, Inches(1.5), Inches(5.5), Inches(4), Inches(0.5),
                 "SAILPOINT  |  DATA PLATFORM", font_size=16, font_color=MID_GRAY,
                 bold=True, alignment=PP_ALIGN.LEFT)

    # Date
    add_text_box(slide, Inches(9), Inches(5.5), Inches(3.5), Inches(0.5),
                 "March 2026  |  Confidential", font_size=14, font_color=MID_GRAY,
                 bold=False, alignment=PP_ALIGN.RIGHT)

    return slide


def create_section_slide(prs, section_title, section_number=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BLUE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ORANGE)
    add_rect(slide, Inches(0), Inches(7.0), SLIDE_WIDTH, Inches(0.08), TEAL)

    if section_number:
        add_text_box(slide, Inches(1.5), Inches(2.2), Inches(2), Inches(1),
                     section_number, font_size=72, font_color=ORANGE, bold=True)

    add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(1.5),
                 section_title, font_size=40, font_color=WHITE, bold=True)

    return slide


def create_content_slide(prs, title, footer_text="SAILPOINT  |  RAPTOR PROJECT"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Top bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), TEAL)

    # Title bar
    add_rect(slide, Inches(0), Inches(0.06), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.7),
                 title, font_size=26, font_color=WHITE, bold=True)

    # Bottom bar
    add_rect(slide, Inches(0), Inches(7.1), SLIDE_WIDTH, Inches(0.4), DARK_BLUE)
    add_text_box(slide, Inches(0.5), Inches(7.12), Inches(5), Inches(0.35),
                 footer_text, font_size=10, font_color=MID_GRAY, bold=False)

    return slide


def add_card(slide, left, top, width, height, title, body_lines, icon="",
             header_color=TEAL, body_color=None):
    card_bg = body_color or RGBColor(0x1A, 0x2D, 0x5A)
    card = add_shape(slide, left, top, width, height, card_bg, border_color=header_color, border_width=1.5)

    # Header stripe
    add_rect(slide, left, top, width, Inches(0.06), header_color)

    title_text = f"{icon}  {title}" if icon else title
    add_text_box(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.4),
                 title_text, font_size=15, font_color=header_color, bold=True)

    y_offset = top + Inches(0.55)
    for line in body_lines:
        add_text_box(slide, left + Inches(0.15), y_offset, width - Inches(0.3), Inches(0.28),
                     line, font_size=11, font_color=LIGHT_GRAY, bold=False)
        y_offset += Inches(0.25)


def add_metric_card(slide, left, top, width, height, metric_value, metric_label,
                    accent_color=TEAL):
    card = add_shape(slide, left, top, width, height, RGBColor(0x1A, 0x2D, 0x5A),
                     border_color=accent_color, border_width=2)
    add_text_box(slide, left, top + Inches(0.15), width, Inches(0.5),
                 metric_value, font_size=32, font_color=accent_color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.65), width - Inches(0.2), Inches(0.4),
                 metric_label, font_size=11, font_color=LIGHT_GRAY, bold=False,
                 alignment=PP_ALIGN.CENTER)


# ==============================================================================
# PPT 1: EXECUTIVE OVERVIEW
# ==============================================================================
def create_exec_overview():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # --- Slide 1: Title ---
    create_title_slide(prs,
                       "🦅  Raptor",
                       "One Door to Data at Scale  —  Executive Overview")

    # --- Slide 2: The Challenge ---
    slide = create_content_slide(prs, "The Challenge: Why We Need Raptor")

    problems = [
        ("💰", "Unsustainable Cost", "Materializers cost $100K+/month.\nRe-processing millions of records for small changes.", ORANGE),
        ("🕸️", "Growing Complexity", "Each feature adds bespoke pipelines.\nBrittle web of dependencies slows innovation.", YELLOW),
        ("🚧", "Platform Bottleneck", "Data platform is critical path for all teams.\nSmall product changes need deep DE involvement.", RED),
        ("⏰", "Fragile Freshness", "\"Real-time\" = frequent batch runs.\nHours of lag before fresh data in products.", ACCENT_PURPLE),
        ("💣", "High Blast Radius", "Minor updates cascade through millions of rows.\nDays of re-materialization work.", TEAL),
    ]

    for i, (icon, title, desc, color) in enumerate(problems):
        col = i % 5
        left = Inches(0.5) + Inches(col * 2.5)
        add_card(slide, left, Inches(1.4), Inches(2.3), Inches(2.6),
                 title, desc.split("\n"), icon=icon, header_color=color)

    # Key stat
    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(0.5),
                 "⚠️  Current search materialization costs exceed $150,000/month in compute alone",
                 font_size=18, font_color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

    # Quote
    add_shape(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(1.4),
              RGBColor(0x1A, 0x2D, 0x5A), border_color=TEAL, border_width=1)
    add_text_box(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(1.0),
                 "\"Our current approach is hitting a wall — technically, financially, and organizationally.\n"
                 "Aperture is the necessary pivot — a shift from pre-built documents to composable, on-demand data delivery.\"",
                 font_size=14, font_color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

    # --- Slide 3: What is Raptor ---
    slide = create_content_slide(prs, "What is Raptor?")

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8),
                 "A single external API that powers every product experience needing data,\n"
                 "delivered quickly, accurately, and cost-effectively.",
                 font_size=20, font_color=LIGHT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    # Before vs After
    add_card(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.5),
             "BEFORE (Materializer-first)", [
                 "❌  Big pre-built documents per use case",
                 "❌  Expensive backfills & replays",
                 "❌  Multi-day lead time for new fields",
                 "❌  \"Freshness\" via constant re-materialization",
                 "❌  Every change = millions of records rebuilt",
                 "❌  High storage from data duplication",
             ], icon="🔴", header_color=RED)

    add_card(slide, Inches(7), Inches(2.5), Inches(5.5), Inches(3.5),
             "AFTER (Raptor / GraphQL-first)", [
                 "✅  On-demand composition via federated subgraphs",
                 "✅  Real-time by design: fetch fresh state directly",
                 "✅  New fields in hours, not weeks",
                 "✅  Materializers become thin or optional",
                 "✅  Targeted rebuilds, not global reprocessing",
                 "✅  Single source of truth per entity",
             ], icon="🟢", header_color=GREEN)

    # Restaurant analogy
    add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.65),
              RGBColor(0x1A, 0x2D, 0x5A), border_color=ORANGE, border_width=1)
    add_text_box(slide, Inches(1.0), Inches(6.25), Inches(11.3), Inches(0.55),
                 "🍽️  Think of it like a restaurant: Today we cook every item on the menu just in case. "
                 "With Raptor, we prepare ingredients once and cook dishes fresh when ordered.",
                 font_size=13, font_color=LIGHT_ORANGE, bold=False, alignment=PP_ALIGN.CENTER)

    # --- Slide 4: Architecture ---
    slide = create_content_slide(prs, "Architecture: 5 Specialized Data Lanes")

    lanes = [
        ("⚡", "Pointed\nLookups", "ScyllaDB / Aerospike", "p95 ≤ 10ms", "Entity hydration\nin milliseconds", TEAL),
        ("🔗", "Graph\nDatabase", "JanusGraph on Scylla", "p95 < 1s", "Relationship\ntraversals", GREEN),
        ("📊", "Real-Time\nAggregates", "ClickHouse", "p95 < 300ms", "Dashboards &\nanomalies", ORANGE),
        ("🔍", "Search &\nDiscovery", "OpenSearch / ES", "p95 < 200ms", "Full-text search\n& autocomplete", ACCENT_PURPLE),
        ("📚", "Historical\n& Batch", "Snowflake", "Hours", "Reporting,\ncompliance, ML", MID_GRAY),
    ]

    for i, (icon, name, backend, slo, purpose, color) in enumerate(lanes):
        left = Inches(0.5) + Inches(i * 2.5)
        card = add_shape(slide, left, Inches(1.5), Inches(2.3), Inches(3.2),
                         RGBColor(0x1A, 0x2D, 0x5A), border_color=color, border_width=2)
        add_rect(slide, left, Inches(1.5), Inches(2.3), Inches(0.06), color)

        add_text_box(slide, left, Inches(1.65), Inches(2.3), Inches(0.4),
                     icon, font_size=28, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.05), Inches(2.3), Inches(0.55),
                     name, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.65), Inches(2.3), Inches(0.3),
                     backend, font_size=10, font_color=color, bold=False, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(3.0), Inches(2.3), Inches(0.35),
                     slo, font_size=20, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.1), Inches(3.45), Inches(2.1), Inches(0.5),
                     purpose, font_size=10, font_color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

    # GraphQL unifier
    add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.8),
              RGBColor(0x1A, 0x2D, 0x5A), border_color=TEAL, border_width=2)
    add_text_box(slide, Inches(0.5), Inches(5.05), Inches(12.3), Inches(0.7),
                 "🌐  Unified GraphQL Gateway  (Apollo Router + Federation 2)  →  Single External Schema",
                 font_size=18, font_color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    # MCP / AI
    add_text_box(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.5),
                 "🤖  AI-Ready via MCP: Schema & metadata auto-published for AI agents and copilots",
                 font_size=14, font_color=LIGHT_ORANGE, bold=False, alignment=PP_ALIGN.CENTER)

    # --- Slide 5: Target Outcomes ---
    slide = create_content_slide(prs, "Target Outcomes & Business Impact")

    metrics = [
        ("30-50%", "Cost Reduction\nvs. Snowflake baseline", TEAL),
        ("< 2 days", "New Entity\nto GraphQL API", GREEN),
        ("≥ 3", "Materializers\nRetired by mid-2026", ORANGE),
        ("99.99%", "Availability\nTarget", ACCENT_PURPLE),
        ("8/10", "Developer\nSatisfaction", LIGHT_TEAL),
    ]

    for i, (val, label, color) in enumerate(metrics):
        left = Inches(0.5) + Inches(i * 2.5)
        add_metric_card(slide, left, Inches(1.5), Inches(2.3), Inches(1.5), val, label, color)

    # Product impact cards
    impacts = [
        ("🚀", "Faster Delivery", "Product teams no longer need custom\npipelines — days instead of weeks"),
        ("⏱️", "Real-Time Experiences", "Fresh data for JIT access, live\ndashboards, dynamic policy eval"),
        ("💰", "Lower Cost", "Reduce data duplication, fewer\nbreakages, simpler operations"),
        ("🤖", "AI-Ready Foundation", "GraphQL + MCP enables natural language\nquerying and AI copilot integration"),
    ]

    for i, (icon, title, desc) in enumerate(impacts):
        left = Inches(0.5) + Inches(i * 3.1)
        add_card(slide, left, Inches(3.5), Inches(2.9), Inches(1.9),
                 title, desc.split("\n"), icon=icon, header_color=TEAL)

    # Bottom quote
    add_text_box(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
                 "\"Raptor gives product teams a single, flexible, real-time source of truth for every experience that depends on data.\"",
                 font_size=16, font_color=LIGHT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    # --- Slide 6: Thank You ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), TEAL)
    add_rect(slide, Inches(0), Inches(7.0), SLIDE_WIDTH, Inches(0.08), ORANGE)
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
                 "Thank You", font_size=52, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
                 "Questions & Discussion", font_size=24, font_color=TEAL, bold=False, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
                 "📧 dattu.marneni@sailpoint.com  |  📋 DPINTAKE-33  |  💬 #raptor-data-platform",
                 font_size=14, font_color=MID_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

    prs.save("Raptor_Executive_Overview.pptx")
    print("✅  Created: Raptor_Executive_Overview.pptx")


# ==============================================================================
# PPT 2: TECHNICAL DEEP DIVE
# ==============================================================================
def create_tech_deep_dive():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # --- Title ---
    create_title_slide(prs,
                       "🦅  Raptor — Technical Deep Dive",
                       "Federated GraphQL Architecture  |  Data Lanes  |  POC Results")

    # --- Slide 2: GraphQL Federation ---
    slide = create_content_slide(prs, "GraphQL Federation: How It Works")

    # Flow boxes
    steps = [
        ("1", "Publish", "Services publish\nentity topics\nto Kafka", TEAL),
        ("2", "Materialize", "Data hydrated,\ndeduped, shaped\ninto clean entities", GREEN),
        ("3", "Fan Out", "Entities distributed\nto fit-for-purpose\nstorage lanes", ORANGE),
        ("4", "Compose", "GraphQL subgraphs\nfederate into\nsingle supergraph", ACCENT_PURPLE),
        ("5", "Serve", "Apollo Router\nserves unified\nexternal API", LIGHT_TEAL),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        left = Inches(0.4) + Inches(i * 2.55)
        card = add_shape(slide, left, Inches(1.4), Inches(2.3), Inches(2.2),
                         RGBColor(0x1A, 0x2D, 0x5A), border_color=color, border_width=2)
        add_text_box(slide, left, Inches(1.5), Inches(2.3), Inches(0.45),
                     num, font_size=32, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(1.95), Inches(2.3), Inches(0.35),
                     title, font_size=16, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.1), Inches(2.4), Inches(2.1), Inches(0.8),
                     desc, font_size=11, font_color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            add_text_box(slide, left + Inches(2.3), Inches(2.0), Inches(0.25), Inches(0.4),
                         "→", font_size=24, font_color=MID_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    # Federation concepts
    concepts = [
        ("@key", "Entity Identity — defines primary key for federated entities"),
        ("@shareable", "Shared Fields — field can be resolved by multiple subgraphs"),
        ("@external", "External Fields — declares field owned by another subgraph"),
        ("extend type Query", "Extending Root Query from subgraphs"),
    ]

    add_text_box(slide, Inches(0.5), Inches(3.9), Inches(4), Inches(0.4),
                 "Key Federation Directives:", font_size=16, font_color=ORANGE, bold=True)

    for i, (directive, desc) in enumerate(concepts):
        y = Inches(4.4) + Inches(i * 0.4)
        add_text_box(slide, Inches(0.8), y, Inches(2), Inches(0.35),
                     directive, font_size=13, font_color=TEAL, bold=True)
        add_text_box(slide, Inches(3.0), y, Inches(9), Inches(0.35),
                     desc, font_size=12, font_color=LIGHT_GRAY, bold=False)

    # Technical requirements
    add_text_box(slide, Inches(0.5), Inches(6.1), Inches(12), Inches(0.4),
                 "Gateway Overhead: ≤5ms p95  |  Availability: ≥99.99%  |  Scale: 100K RPS  |  "
                 "Cache Invalidation: <1s  |  Schema Rollback: ≤5min",
                 font_size=12, font_color=ORANGE, bold=False, alignment=PP_ALIGN.CENTER)

    # --- Slide 3: Subgraphs ---
    slide = create_content_slide(prs, "Subgraph Architecture: aperture-core & aperture-graph")

    # Core subgraph
    add_card(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.0),
             "aperture-core  (Pointed Lookups)", [
                 "Backend: ScyllaDB (key-value lookups)",
                 "",
                 "Owns core entity types:",
                 "  • Identity, Entitlement, Source",
                 "  • AccessProfile, Role, Account",
                 "",
                 "All base properties:",
                 "  • name, email, department, etc.",
                 "  • Direct relationships: manager, accounts",
                 "  • entitlementAssignments",
                 "",
                 "SLO: p95 ≤ 10ms  |  Freshness: < 60s",
                 "Scale: ≥ 100M identities, ≥ 5B records",
                 "",
                 "Cost target: 30-50% reduction vs current",
             ], icon="⚡", header_color=TEAL)

    add_card(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(5.0),
             "aperture-graph  (Graph Traversals)", [
                 "Backend: JanusGraph on Scylla",
                 "",
                 "Owns graph-specific fields:",
                 "  • effectiveAccess, identitiesWithAccess",
                 "  • traverseIdentities, traverseEntitlements",
                 "  • searchIdentities (nested filtering)",
                 "  • Access path analysis",
                 "",
                 "Relationship chain:",
                 "  identity → role → access_profile →",
                 "  entitlement → resource",
                 "",
                 "SLO: p95 < 1s (N-hop), < 100ms (10-hop)",
                 "Scale: ≥ 5B edges across ≥ 100M nodes",
                 "",
                 "Enables: policy eval, JIT access, SoD",
             ], icon="🔗", header_color=GREEN)

    # Cross-subgraph note
    add_shape(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
              RGBColor(0x1A, 0x2D, 0x5A), border_color=ORANGE, border_width=1)
    add_text_box(slide, Inches(0.5), Inches(6.52), Inches(12.3), Inches(0.45),
                 "🔄  Cross-Subgraph: Router automatically stitches core + graph data. "
                 "One query fetches identity name (core) + effective access (graph) seamlessly.",
                 font_size=13, font_color=LIGHT_ORANGE, bold=False, alignment=PP_ALIGN.CENTER)

    # --- Slide 4: Use Cases ---
    slide = create_content_slide(prs, "Key Use Cases Powered by Raptor")

    use_cases = [
        ("🔍", "Identity Quick View", "Single query: name, email, org +\nall effective entitlements.\nNo separate lookups needed.", TEAL),
        ("🔎", "Search with Context", "Search identities + enrich with\nprofile info + real-time activity\ncounts in one call.", GREEN),
        ("🛡️", "Access Impact Analysis", "\"If I revoke this entitlement,\nwho is affected and how have\nthey used it?\"", ORANGE),
        ("⚖️", "Policy Engine Context", "On-demand identity context for\nSoD, JIT access, continuous\ncompliance checks.", ACCENT_PURPLE),
        ("🕸️", "Identity Graph", "Visual exploration of access\nrelationships, path tracing,\nprivilege exposure analysis.", LIGHT_TEAL),
        ("🔄", "Role Change Propagation", "Graph-scoped diffs for role\nchanges → targeted provision/\ndeprovision actions.", YELLOW),
        ("📊", "Real-Time Dashboards", "Live login counts, access\nrequests, anomaly detection\nwithout Snowflake scans.", ORANGE),
        ("🔐", "Simplified Materialization", "Pointed lookups replace stateful\nstreaming joins. Identity docs\nbecome stateless enrichment.", TEAL),
    ]

    for i, (icon, title, desc, color) in enumerate(use_cases):
        row = i // 4
        col = i % 4
        left = Inches(0.4) + Inches(col * 3.15)
        top = Inches(1.3) + Inches(row * 2.7)
        add_card(slide, left, top, Inches(3.0), Inches(2.3),
                 title, desc.split("\n"), icon=icon, header_color=color)

    # --- Slide 5: Technology Stack ---
    slide = create_content_slide(prs, "Technology Decisions & Trade-offs")

    decisions = [
        ("GraphQL Gateway", "Apollo Router + Federation 2", "Excellent",
         "High-perf Rust gateway, mature federation,\npreview/compose, MCP integration", TEAL),
        ("Low-Latency Store", "ScyllaDB (self-hosted)", "Excellent",
         "Ultra-low latency, cost-efficient,\nFedRAMP-ready, same cluster as graph", GREEN),
        ("Graph Database", "JanusGraph + Scylla", "Excellent",
         "Mature open source, scalable,\nleverages same Scylla cluster", ORANGE),
        ("Real-Time Analytics", "ClickHouse", "Excellent",
         "Extreme read speed, columnar storage,\nmature ecosystem, Kafka connectors", ACCENT_PURPLE),
        ("Alternative: Lookups", "Aerospike SaaS", "Excellent",
         "Very low latency, built-in HA,\nbeing evaluated as alternative", LIGHT_TEAL),
    ]

    for i, (lane, choice, fit, notes, color) in enumerate(decisions):
        top = Inches(1.3) + Inches(i * 1.1)
        add_rect(slide, Inches(0.5), top, Inches(0.08), Inches(0.9), color)
        add_text_box(slide, Inches(0.8), top + Inches(0.05), Inches(2.5), Inches(0.35),
                     lane, font_size=15, font_color=WHITE, bold=True)
        add_text_box(slide, Inches(3.5), top + Inches(0.05), Inches(3.5), Inches(0.35),
                     choice, font_size=15, font_color=color, bold=True)
        add_text_box(slide, Inches(7.2), top + Inches(0.05), Inches(1.2), Inches(0.35),
                     f"✅ {fit}", font_size=12, font_color=GREEN, bold=False)
        add_text_box(slide, Inches(8.5), top + Inches(0.05), Inches(4.5), Inches(0.8),
                     notes, font_size=11, font_color=LIGHT_GRAY, bold=False)

    # Repos
    add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.3),
                 "📂 Repos: aperture-core (ScyllaDB) | aperture-graph (JanusGraph) | "
                 "aperture-search (ES) | pointed-lookups-dip | gitops-k8s (apollo-poc)",
                 font_size=11, font_color=MID_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

    # --- Slide 6: Thank You ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), TEAL)
    add_rect(slide, Inches(0), Inches(7.0), SLIDE_WIDTH, Inches(0.08), ORANGE)
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
                 "Deep Dive Complete", font_size=48, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
                 "Let's Build the Future of Data at SailPoint", font_size=24, font_color=TEAL,
                 bold=False, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
                 "📧 dattu.marneni@sailpoint.com  |  📋 DPINTAKE-33  |  💬 #raptor-data-platform",
                 font_size=14, font_color=MID_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

    prs.save("Raptor_Technical_Deep_Dive.pptx")
    print("✅  Created: Raptor_Technical_Deep_Dive.pptx")


# ==============================================================================
# PPT 3: ROADMAP & STATUS
# ==============================================================================
def create_roadmap_status():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # --- Title ---
    create_title_slide(prs,
                       "🦅  Raptor — Roadmap & Status",
                       "Phased Delivery  |  DACI Status  |  Current Progress  |  Next Steps")

    # --- Slide 2: Phased Roadmap ---
    slide = create_content_slide(prs, "Phased Delivery Roadmap")

    phases = [
        ("Phase I", "Prove It Works", "Q4'25 - Q1'26", GREEN,
         ["Validate each lane via POCs", "Joint observability dashboards",
          "Runbooks & cost models", "ALT decision gate: continue/pivot/pause"]),
        ("Phase II", "Golden Paths", "Q2'26", YELLOW,
         ["Documented Kafka→Store→GraphQL flows", "Schema & SLO templates",
          "First product integration", "MySailPoint or Access Center"]),
        ("Phase III", "Self-Service", "Q3'26", TEAL,
         ["Prebuilt templates", "Onboarding guides & cost dashboards",
          "Internal workshops", "Scale early adopters"]),
        ("Phase IV", "Scale & Harden", "Q4'26+", ACCENT_PURPLE,
         ["External GraphQL exposure", "FedRAMP readiness",
          "Cost optimization", "SLO enforcement"]),
    ]

    for i, (phase, name, timeline, color, items) in enumerate(phases):
        left = Inches(0.4) + Inches(i * 3.15)
        card = add_shape(slide, left, Inches(1.3), Inches(3.0), Inches(4.8),
                         RGBColor(0x1A, 0x2D, 0x5A), border_color=color, border_width=2)
        add_rect(slide, left, Inches(1.3), Inches(3.0), Inches(0.06), color)

        add_text_box(slide, left, Inches(1.45), Inches(3.0), Inches(0.35),
                     phase, font_size=14, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(1.8), Inches(3.0), Inches(0.4),
                     name, font_size=18, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.2), Inches(3.0), Inches(0.3),
                     timeline, font_size=12, font_color=color, bold=False, alignment=PP_ALIGN.CENTER)

        for j, item in enumerate(items):
            add_text_box(slide, left + Inches(0.15), Inches(2.65) + Inches(j * 0.38),
                         Inches(2.7), Inches(0.35),
                         f"• {item}", font_size=11, font_color=LIGHT_GRAY, bold=False)

    # Current marker
    add_text_box(slide, Inches(0.4), Inches(6.3), Inches(3.0), Inches(0.4),
                 "◀ WE ARE HERE", font_size=14, font_color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    # --- Slide 3: DACI Status ---
    slide = create_content_slide(prs, "DACI Decision Status & POC Progress")

    dacis = [
        ("🌐", "GraphQL Layer", "Apollo Router + Federation 2", "📝 Draft",
         "SAASD8NG-4502", "Stand up router,\nfirst subgraph integrated", TEAL),
        ("⚡", "Pointed Lookups", "ScyllaDB (self-hosted + Cloud)", "🔄 In Progress",
         "SAASD8NG-4501", "p95 ≤ 10ms @ 5K RPS,\ncost ↓ 30%", GREEN),
        ("🔗", "Graph Database", "JanusGraph on Scylla", "🔄 In Progress",
         "SAASD8NG-4273", "p95 < 1s for N-hop,\n5B edges scale", ORANGE),
        ("📊", "Real-Time Analytics", "ClickHouse (self-hosted/SaaS)", "📝 Draft",
         "SAASD8NG-4972", "p95 < 300ms,\ncost ↓ 40% vs SF", ACCENT_PURPLE),
        ("🔍", "Search Layer", "Aperture Search subgraph", "📝 Draft",
         "—", "Search 2.0 strategy,\ncontract compatibility", LIGHT_TEAL),
    ]

    for i, (icon, name, tech, status, epic, criteria, color) in enumerate(dacis):
        top = Inches(1.3) + Inches(i * 1.1)
        add_rect(slide, Inches(0.5), top, Inches(0.08), Inches(0.95), color)

        add_text_box(slide, Inches(0.7), top + Inches(0.05), Inches(0.4), Inches(0.3),
                     icon, font_size=18, font_color=color, bold=True)
        add_text_box(slide, Inches(1.2), top + Inches(0.05), Inches(2.0), Inches(0.3),
                     name, font_size=15, font_color=WHITE, bold=True)
        add_text_box(slide, Inches(3.4), top + Inches(0.05), Inches(2.8), Inches(0.3),
                     tech, font_size=12, font_color=LIGHT_GRAY, bold=False)
        status_color = GREEN if "Progress" in status else YELLOW
        add_text_box(slide, Inches(6.4), top + Inches(0.05), Inches(1.5), Inches(0.3),
                     status, font_size=12, font_color=status_color, bold=True)
        add_text_box(slide, Inches(8.0), top + Inches(0.05), Inches(1.8), Inches(0.3),
                     epic, font_size=11, font_color=TEAL, bold=False)
        add_text_box(slide, Inches(10.0), top + Inches(0.05), Inches(3.0), Inches(0.85),
                     criteria, font_size=10, font_color=LIGHT_GRAY, bold=False)

    # Action tracking
    add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.3),
                 "Driver: Derrick Mink (departed) → Ownership gap to be addressed  |  "
                 "Approvers: Dan Sparks, Fuad Rashid  |  INIT-2073: In Progress",
                 font_size=11, font_color=ORANGE, bold=False, alignment=PP_ALIGN.CENTER)

    # --- Slide 4: Risks & Actions ---
    slide = create_content_slide(prs, "Risks, Mitigations & Immediate Actions")

    risks = [
        ("⚠️", "Driver Departed", "Derrick Mink left — ownership gap\nacross all DACI lanes", "Assign new drivers per lane\nimmediately", RED),
        ("🚧", "DevOps Bottleneck", "Infrastructure intake can block\nPOC velocity", "Pre-stub stories, align early\nwith SRE teams", ORANGE),
        ("🏗️", "Infra Complexity", "Self-hosted DBs require ops\nexpertise", "Evaluate SaaS alternatives\nin parallel", YELLOW),
        ("👥", "Team Capacity", "Limited bandwidth for\nmulti-lane POCs", "Engage evolv/tiger teams\nfor environment setup", TEAL),
    ]

    for i, (icon, title, risk, mitigation, color) in enumerate(risks):
        col = i % 4
        left = Inches(0.4) + Inches(col * 3.15)
        add_card(slide, left, Inches(1.3), Inches(3.0), Inches(2.3),
                 title, (risk + "\n\n✅ " + mitigation).split("\n"),
                 icon=icon, header_color=color)

    # Next steps
    actions = [
        ("1️⃣", "Assess current POC state — what's deployed & running"),
        ("2️⃣", "Clarify ownership — new drivers for each DACI lane"),
        ("3️⃣", "Address DPINTAKE-33 — delivery date for Identity Graph"),
        ("4️⃣", "Unblock backlog — assign subgraph stories 5980/5981/5982"),
        ("5️⃣", "Rename to Raptor — update Confluence, Jira, Slack"),
        ("6️⃣", "Establish weekly cadence — \"Raptor Updates\" posts"),
        ("7️⃣", "Schedule kickoff sync — Forge, Keel, DIPO, Double Zero"),
    ]

    add_text_box(slide, Inches(0.5), Inches(3.9), Inches(5), Inches(0.4),
                 "⚡ Immediate Next Steps:", font_size=18, font_color=ORANGE, bold=True)

    for i, (num, action) in enumerate(actions):
        col = i // 4
        row = i % 4
        left = Inches(0.5) + Inches(col * 6.3)
        top = Inches(4.4) + Inches(row * 0.5)
        add_text_box(slide, left, top, Inches(6.0), Inches(0.45),
                     f"{num}  {action}", font_size=13, font_color=LIGHT_GRAY, bold=False)

    # --- Slide 5: Thank You ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), TEAL)
    add_rect(slide, Inches(0), Inches(7.0), SLIDE_WIDTH, Inches(0.08), ORANGE)
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
                 "Let's Move Forward", font_size=48, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
                 "Raptor: From POC to Production", font_size=24, font_color=TEAL,
                 bold=False, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
                 "📧 dattu.marneni@sailpoint.com  |  📋 DPINTAKE-33  |  💬 #raptor-data-platform",
                 font_size=14, font_color=MID_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

    prs.save("Raptor_Roadmap_and_Status.pptx")
    print("✅  Created: Raptor_Roadmap_and_Status.pptx")


if __name__ == "__main__":
    create_exec_overview()
    create_tech_deep_dive()
    create_roadmap_status()
    print("\n🎉  All 3 presentations created successfully!")
